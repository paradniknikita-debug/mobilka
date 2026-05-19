"""
Презентация 15 слайдов в стиле исходного PPTX (макеты Slide 1–10 master).

Запуск из корня проекта:
  py -3 scripts/build_presentation_from_template.py

Результат:
  Informacionnyj-resurs-dlya-verifikacii-dannyh-geolokacii-obuektov-elektroenergetiki-15slides.pptx

На слайдах с картинкой слева (макет 6/8/10) — замените изображение:
  ПКМ по рисунку → «Изменить рисунок» → ваш скриншот.
"""
from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "Informacionnyj-resurs-dlya-verifikacii-dannyh-geolokacii-obuektov-elektroenergetiki.pptx"
OUT = ROOT / "Informacionnyj-resurs-dlya-verifikacii-dannyh-geolokacii-obuektov-elektroenergetiki-15slides.pptx"

def _set_text(shape, text: str) -> None:
    if not hasattr(shape, "text"):
        return
    shape.text = text


def _set_pair(slide, label_idx: int, value_idx: int, label: str, value: str) -> None:
    _set_text(slide.shapes[label_idx], label)
    _set_text(slide.shapes[value_idx], value)


def _fill_bullets_pic_right(slide, title: str, bullets: list[str], screenshot_hint: str) -> None:
    """Макет Slide 6 master: картинка слева, справа заголовок и маркеры."""
    shapes = slide.shapes
    _set_text(shapes[1], title)  # Text 0 is index 1 after Image 0
    bullet_indices = [3, 5, 7, 9]
    for i, idx in enumerate(bullet_indices):
        if i < len(bullets):
            _set_text(shapes[idx], bullets[i])
        else:
            _set_text(shapes[idx], "")
    # Подсказка в последний видимый блок или в неиспользуемый
    for j in range(len(bullets), len(bullet_indices)):
        _set_text(shapes[bullet_indices[j]], screenshot_hint if j == len(bullets) and screenshot_hint else "")


def _fill_pic_left_blocks(slide, title: str, blocks: list[tuple[str, str]], screenshot_hint: str) -> None:
    """Макет Slide 8 master: картинка слева, блоки заголовок+текст."""
    shapes = slide.shapes
    _set_text(shapes[1], title)
    pairs = [(2, 3), (4, 5), (6, 7), (8, 9)]
    for i, (li, vi) in enumerate(pairs):
        if i < len(blocks):
            _set_pair(slide, li, vi, blocks[i][0], blocks[i][1])
        else:
            _set_pair(slide, li, vi, "", "")
    if screenshot_hint and len(blocks) < 4:
        li, vi = pairs[len(blocks)]
        _set_pair(slide, li, vi, "Скриншот", screenshot_hint)


def _fill_numbered(slide, title: str, goal: str, tasks: list[str]) -> None:
    shapes = slide.shapes
    _set_text(shapes[0], title)
    _set_text(shapes[1], goal)
    _set_text(shapes[2], "Задачи:")
  # pairs: num at 3,5,7... text at 5,8,11...
    num_indices = [3, 6, 9, 12, 15, 18]
    text_indices = [5, 8, 11, 14, 17, 20]
    for i in range(min(len(tasks), len(num_indices))):
        _set_text(shapes[num_indices[i]], f"{i + 1:02d}")
        _set_text(shapes[text_indices[i]], tasks[i])
    for j in range(len(tasks), len(num_indices)):
        _set_text(shapes[num_indices[j]], "")
        _set_text(shapes[text_indices[j]], "")


def _fill_architecture(slide) -> None:
    pairs = [
        (2, 3, "Клиент", "Веб-приложение (Angular)"),
        (5, 6, "Точка входа", "Nginx, HTTPS"),
        (8, 9, "Сервер", "Python / FastAPI, REST API"),
        (11, 12, "Данные", "PostgreSQL, геокоординаты опор и линий"),
        (14, 15, "Кэш", "Redis"),
        (17, 18, "Медиа", "Фото и файлы с обходов (S3 / MinIO)"),
        (20, 21, "Развёртывание", "Docker Compose"),
    ]
    _set_text(slide.shapes[0], "Архитектура решения")
    for li, vi, lab, val in pairs:
        _set_pair(slide, li, vi, lab, val)


def _fill_table_roles(slide) -> None:
    _set_text(slide.shapes[0], "Роли пользователей")
    _set_text(slide.shapes[3], "Роль")
    _set_text(slide.shapes[4], "Назначение")
    rows = [
        (6, 7, "Администратор", "Учётные записи, мониторинг сервисов"),
        (9, 10, "Паспортист", "Паспорта, справочники, выгрузки CIM и отчёты"),
        (12, 13, "Инженер-обходчик", "Карта, редактирование объектов, вложения"),
        (15, 16, "Безопасность", "Вход по логину и паролю (JWT)"),
        (18, 19, "", ""),
        (21, 22, "", ""),
    ]
    for li, vi, lab, val in rows:
        if lab:
            _set_pair(slide, li, vi, lab, val)
        else:
            _set_text(slide.shapes[li], "")
            _set_text(slide.shapes[vi], "")


def _fill_cim(slide) -> None:
    _set_text(slide.shapes[0], "Интеграция CIM")
    _set_pair(slide, 2, 3, "Экспорт", "Модель сети в CIM XML (вся сеть или одна ЛЭП)")
    _set_pair(slide, 5, 6, "Импорт", "Предпросмотр и применение изменений")
    _set_pair(slide, 8, 9, "Стандарт", "IEC 61970 — обмен с внешними системами")
    _set_text(slide.shapes[10], "Замените три рисунка: экран CIM, фрагмент XML, схема обмена")


def _fill_results(slide) -> None:
    _set_text(slide.shapes[1], "Результаты")
    bullets = [
        "Веб-система учёта ЛЭП, опор, подстанций и оборудования",
        "Карта, паспортизация, журнал изменений, администрирование",
        "Экспорт и импорт CIM XML, контроль согласованности данных",
        "Готовность к развёртыванию в Docker",
    ]
    text_idx = [3, 5, 7, 9]
    for i, idx in enumerate(text_idx):
        if i < len(bullets):
            _set_text(slide.shapes[idx], bullets[i])


def _fill_title(slide) -> None:
    shapes = slide.shapes
    _set_text(
        shapes[1],
        "Информационный ресурс для верификации данных геолокации объектов электроэнергетики",
    )
    _set_text(
        shapes[2],
        "Веб-система учёта и визуализации объектов электросетевого хозяйства",
    )
    _set_text(
        shapes[3],
        "Исполнитель: Парадник Никита Владимирович\nгр. 10703222\n"
        "Руководитель: Гутич Ирина Ивановна",
    )
    _set_text(shapes[4], "БНТУ, 2026")


def _duplicate_slide(prs: Presentation, source_index: int):
    """Копия слайда со всеми фигурами (макеты в файле без плейсхолдеров)."""
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in list(dest.shapes):
        el = shape.element
        el.getparent().remove(el)
    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return dest


def build() -> None:
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    shutil.copy(SRC, OUT)
    prs = Presentation(OUT)

    # --- Обновляем 10 исходных слайдов ---
    _fill_title(prs.slides[0])

    _fill_numbered(
        prs.slides[1],
        "Цель и задачи",
        "Цель: разработка веб-системы учёта и визуализации объектов электросетевого хозяйства с картой и обменом по стандарту CIM.",
        [
            "Анализ предметной области и требований",
            "Проектирование архитектуры и базы данных",
            "Реализация серверной логики и веб-интерфейса",
            "Картографический модуль и учёт объектов",
            "Модуль CIM и паспортизация",
            "Развёртывание и тестирование",
        ],
    )

    _fill_architecture(prs.slides[2])
    _fill_table_roles(prs.slides[3])

    # Бывший Frontend-стек → объекты учёта (табличный макет)
    s5 = prs.slides[4]
    _set_text(s5.shapes[0], "Объекты учёта")
    _set_text(s5.shapes[3], "Объект")
    _set_text(s5.shapes[4], "Содержание")
    obj_rows = [
        (6, 7, "ЛЭП", "Участки, пролёты, класс напряжения"),
        (9, 10, "Опоры", "Координаты, оборудование, медиа"),
        (12, 13, "Подстанции", "Привязка линий, паспортные данные"),
        (15, 16, "Оборудование", "Справочник марок, привязка к опорам"),
    ]
    for li, vi, lab, val in obj_rows:
        _set_pair(s5, li, vi, lab, val)

    _fill_bullets_pic_right(
        prs.slides[5],
        "Рабочее место: карта",
        [
            "Интерактивная карта и дерево объектов",
            "Создание и редактирование ЛЭП, опор, подстанций",
            "Карточки объектов и контекстное меню",
            "Вставьте скриншот полного экрана карты",
        ],
        "Скриншот: /map",
    )

    _fill_cim(prs.slides[6])

    _fill_pic_left_blocks(
        prs.slides[7],
        "Медиавложения с обхода",
        [
            ("Типы", "Фото, схемы, голосовые заметки, видео"),
            ("Хранение", "Облачное или локальное хранилище"),
            ("Связь", "Прикрепление к карточке опоры и оборудования"),
        ],
        "Скриншот: вложения в карточке опоры",
    )

    # Бывшее развёртывание → паспортизация (три колонки с картинками — заменить на скрины вкладок)
    s9 = prs.slides[8]
    _set_text(s9.shapes[0], "Паспортизация")
    _set_pair(s9, 4, 5, "Отчёты", "Дефекты, сводки, журнал обходов")
    _set_pair(s9, 9, 10, "Паспорта", "PDF, Word, Excel по объекту")
    _set_pair(s9, 14, 15, "Справочники", "Марки оборудования и проводов")
    _set_text(s9.shapes[19], "Замените 4 рисунка на скрины раздела «Паспортизация»")

    _fill_results(prs.slides[9])

    # --- Добавляем 5 слайдов (11–15) копированием оформленных слайдов ---
    s11 = _duplicate_slide(prs, 7)
    _fill_pic_left_blocks(
        s11,
        "Карточка опоры",
        [
            ("Атрибуты", "Тип, материал, координаты, дефекты"),
            ("Оборудование", "Список на опоре"),
            ("Медиа", "Фото и файлы с обхода"),
        ],
        "Скриншот: диалог опоры",
    )

    s12 = _duplicate_slide(prs, 5)
    _fill_bullets_pic_right(
        s12,
        "ЛЭП и топология",
        [
            "Участки и пролёты между опорами",
            "Автосборка топологии по цепочке опор",
            "Отпайки и сегменты линии",
            "Скриншот: карточка участка или пересборка",
        ],
        "",
    )

    s13 = _duplicate_slide(prs, 3)
    _set_text(s13.shapes[0], "Качество данных")
    _set_text(s13.shapes[3], "Механизм")
    _set_text(s13.shapes[4], "Описание")
    qual = [
        (6, 7, "Валидация", "Согласованность класса напряжения ЛЭП и объектов"),
        (9, 10, "Журнал", "История изменений по пользователям"),
        (12, 13, "Контроль", "Выявление разрывов топологии"),
    ]
    for li, vi, lab, val in qual:
        _set_pair(s13, li, vi, lab, val)

    s14 = _duplicate_slide(prs, 5)
    _fill_bullets_pic_right(
        s14,
        "Администрирование",
        [
            "Список пользователей и роли",
            "Сводка по БД и сервисам",
            "Скриншот: раздел «Администрирование»",
        ],
        "",
    )

    s15 = _duplicate_slide(prs, 0)
    _fill_title(s15)
    _set_text(s15.shapes[2], "Спасибо за внимание!")
    _set_text(s15.shapes[3], "Готов ответить на вопросы")
    _set_text(s15.shapes[4], "БНТУ, 2026")

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print("Замените изображения (Image): ПКМ → Изменить рисунок.")


if __name__ == "__main__":
    build()
