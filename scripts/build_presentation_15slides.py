"""
Сборка презентации (15 слайдов) с зонами под скриншоты.
Запуск из корня проекта:
  py -3 scripts/build_presentation_15slides.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Informacionnyj-resurs-dlya-verifikacii-dannyh-geolokacii-obuektov-elektroenergetiki-15slides.pptx"

# Палитра (энергетика / корпоратив)
C_BG = RGBColor(0xF4, 0xF7, 0xFB)
C_PRIMARY = RGBColor(0x0D, 0x47, 0xA1)
C_ACCENT = RGBColor(0x00, 0x96, 0x88)
C_TEXT = RGBColor(0x1A, 0x23, 0x32)
C_MUTED = RGBColor(0x5C, 0x6B, 0x7A)
C_PLACEHOLDER = RGBColor(0xE8, 0xED, 0xF3)
C_PLACEHOLDER_BORDER = RGBColor(0x90, 0xA4, 0xAE)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.55))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = C_MUTED


def _add_bullets(slide, items: list[str], left=0.7, top=1.75, width=5.8, height=5.0, font_pt=18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_pt)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(8)


def _add_screenshot_zone(
    slide,
    hint: str,
    *,
    left=6.2,
    top=1.55,
    width=6.5,
    height=5.35,
) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PLACEHOLDER
    shape.line.color.rgb = C_PLACEHOLDER_BORDER
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"📷 СКРИНШОТ\n\n{hint}"
    p.font.size = Pt(14)
    p.font.color.rgb = C_MUTED
    p.font.bold = False


def _add_accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Титул
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.2))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Информационный ресурс для верификации данных геолокации объектов электроэнергетики"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p2 = tf.add_paragraph()
    p2.text = "Веб-система учёта ЛЭП, опор, подстанций и оборудования с картой и обменом по CIM"
    p2.font.size = Pt(20)
    p2.font.color.rgb = C_MUTED
    p2.space_before = Pt(16)
    meta = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11), Inches(1.2))
    mtf = meta.text_frame
    mp = mtf.paragraphs[0]
    mp.text = "Дипломный проект · БНТУ · 2026\nИсполнитель: Парадник Н. В. · гр. 10703222"
    mp.font.size = Pt(16)
    mp.font.color.rgb = C_TEXT
    _add_screenshot_zone(
        s,
        "Опционально: логотип БНТУ или общий вид карты (уменьшенный)",
        left=9.0,
        top=1.2,
        width=3.6,
        height=2.4,
    )

    # 2 — Проблема
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Актуальность и проблема", "Почему нужна цифровизация учёта сетевых объектов")
    _add_bullets(
        s,
        [
            "Данные об опорах и оборудовании часто разрознены (таблицы, бумага).",
            "Нет наглядной привязки объектов к местности на карте.",
            "Сложно оперативно обновлять информацию после обходов и ремонтов.",
            "Интеграция с отраслевыми системами требует единой модели (CIM).",
        ],
    )
    _add_screenshot_zone(s, "Схема «было → стало» или фрагмент бумажного паспорта / Excel (для контраста)")

    # 3 — Цель и задачи
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Цель и задачи проекта")
    _add_bullets(
        s,
        [
            "Цель: веб-система учёта и визуализации объектов электросетевого хозяйства.",
            "Единая база: ЛЭП, опоры, пролёты, подстанции, оборудование.",
            "Карта с геопривязкой и карточки объектов.",
            "Обмен данными в формате CIM XML (стандарты IEC 61970).",
            "Паспортизация, отчёты, контроль качества данных.",
        ],
        width=12.0,
    )

    # 4 — Объекты учёта
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Что учитывает система", "Иерархия объектов на карте и в базе")
    _add_bullets(
        s,
        [
            "Линия электропередачи → участки и пролёты → опоры.",
            "Подстанции и привязка линий к ПС.",
            "Оборудование на опорах (выключатели, разъединители и др.).",
            "Медиавложения: фото, схемы, голосовые заметки с обхода.",
        ],
    )
    _add_screenshot_zone(s, "Дерево объектов в боковой панели (sidebar) с раскрытой ЛЭП")

    # 5 — Архитектура (без файлов)
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Архитектура решения", "Клиент · сервер · база данных · хранилище файлов")
    _add_bullets(
        s,
        [
            "Веб-клиент (Angular) — карта, формы, отчёты.",
            "Серверное приложение (Python / FastAPI) — бизнес-логика и API.",
            "PostgreSQL + геоданные — учёт объектов и координат.",
            "Хранилище вложений (S3 / MinIO или диск).",
            "Развёртывание: Docker, Nginx, HTTPS.",
        ],
        width=5.5,
    )
    _add_screenshot_zone(
        s,
        "Блок-схема: Браузер → Nginx → Backend → БД / MinIO\n(можно нарисовать в draw.io или вставить из ПЗ, упрощённую)",
        left=6.0,
        top=1.5,
        width=6.8,
        height=5.2,
    )

    # 6 — Роли
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Роли пользователей", "Три уровня доступа")
    _add_bullets(
        s,
        [
            "Администратор — пользователи, мониторинг сервисов, полный доступ.",
            "Паспортист — паспорта, справочники, выгрузки CIM и отчёты.",
            "Инженер-обходчик — карта, редактирование объектов, журнал, вложения.",
        ],
        width=5.5,
    )
    _add_screenshot_zone(s, "Раздел «Администрирование»: таблица пользователей и сводка")

    # 7 — Вход и карта
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Рабочее место: карта", "Главный экран системы")
    _add_bullets(
        s,
        [
            "Интерактивная карта (подложка OSM / офлайн-тайлы).",
            "ЛЭП, опоры, подстанции, пролёты — отдельные слои.",
            "Дерево объектов слева, контекстное меню на карте.",
            "Создание и редактирование объектов в диалогах.",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "Полный экран «Карта»: линия, опоры, подстанции, открытый sidebar")

    # 8 — Карточка опоры
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Карточка объекта", "Опора: атрибуты, оборудование, вложения")
    _add_bullets(
        s,
        [
            "Паспортные поля опоры и координаты.",
            "Список оборудования на опоре.",
            "Фото и медиа с обхода, комментарии.",
            "Дефекты и критичность (при необходимости).",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "Диалог карточки опоры: вкладки / оборудование / вложения (фото)")

    # 9 — ЛЭП и топология
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "ЛЭП и электрическая модель", "Участки, пролёты, отпайки")
    _add_bullets(
        s,
        [
            "Редактирование параметров линии (напряжение, диспетчерское имя).",
            "Участки ACLineSegment, пролёты между опорами.",
            "Автосборка топологии по последовательности опор.",
            "Карточка участка с длиной по пролётам.",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "Диалог линии или участка / «Пересборка топологии» / карточка сегмента")

    # 10 — CIM
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Интеграция CIM", "Обмен с внешними системами")
    _add_bullets(
        s,
        [
            "Экспорт модели сети в CIM XML (вся сеть или одна ЛЭП).",
            "Импорт и предпросмотр изменений (diff).",
            "Настройки выгрузки: GPS, оборудование, электрическая модель.",
            "Соответствие отраслевым стандартам IEC 61970.",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "Вкладка «CIM / 552»: экспорт, настройки или предпросмотр импорта")

    # 11 — Паспортизация
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Паспортизация", "Отчёты, техпаспорта, справочники")
    _add_bullets(
        s,
        [
            "Отчёты по дефектам, линии, обходам.",
            "Технические паспорта — формирование и выгрузка PDF / Word / Excel.",
            "Справочник марок оборудования и проводов.",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "Раздел «Паспортизация»: вкладки Отчёты / Паспорта / Справочники")

    # 12 — Валидация и журнал
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Качество данных", "Контроль согласованности и аудит")
    _add_bullets(
        s,
        [
            "Проверка согласованности класса напряжения ЛЭП, опор, каталога, сегментов.",
            "Журнал изменений: кто и когда менял объекты.",
            "Журнал несоответствий топологии (для администратора).",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "«Журнал» или сообщение об ошибке валидации при сохранении")

    # 13 — Мобильный клиент (опционально)
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Полевой доступ", "Мобильный клиент (Flutter)")
    _add_bullets(
        s,
        [
            "Тот же сервер и учётные записи.",
            "Карта и карточки объектов в полевых условиях.",
            "Синхронизация данных при появлении связи.",
        ],
        width=5.4,
    )
    _add_screenshot_zone(s, "Flutter Web или мобильный эмулятор: карта / карточка опоры\n(если показываете на защите)")

    # 14 — Результаты
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_BG)
    _add_accent_bar(s)
    _add_title(s, "Результаты работы")
    _add_bullets(
        s,
        [
            "Реализован программный комплекс: сервер + веб-интерфейс.",
            "Учёт и визуализация объектов сети на карте.",
            "CIM-экспорт/импорт, паспортизация, администрирование.",
            "Контейнерное развёртывание, готовность к интеграции.",
        ],
        width=12.0,
        top=1.7,
    )

    # 15 — Спасибо
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, C_PRIMARY)
    _add_accent_bar(s)
    t = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Готов ответить на вопросы"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
